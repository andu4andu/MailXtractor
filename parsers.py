import csv
import io
import logging
import math
import os
import shutil
import tempfile
import warnings
import zipfile
import openpyxl

warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl")
import pdfplumber
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pytesseract
from pdf2image import convert_from_path

import json as _json
import os as _os
with open("config.json", encoding="utf-8") as _f:
    _config = _json.load(_f)
pytesseract.pytesseract.tesseract_cmd = _os.environ.get("TESSERACT_PATH", _config["tesseract_path"])
_SKIP_SHEETS = _config["skip_sheets"]

logger = logging.getLogger(__name__)

IMAGE_EXTENSIONS = {".png", ".jpeg", ".jpg", ".bmp", ".tiff"}
PDF_EXTENSIONS = {".pdf"}
SPREADSHEET_EXTENSIONS = {".xlsx", ".xls", ".csv"}
PRESENTATION_EXTENSIONS = {".ppt", ".pptx"}
DOCUMENT_EXTENSIONS = {".doc", ".docx", ".txt"}
ZIP_EXTENSIONS = {".zip", ".rar", ".7z"}


def _join_cell(cell: str) -> str:
    parts = cell.split("\n")
    result = parts[0]
    for part in parts[1:]:
        stripped = part.strip()
        is_fragment = (
            result and result[-1].isalpha() and
            stripped and stripped[0].islower() and
            stripped.isalpha() and
            len(stripped) <= 5
        )
        result += stripped if is_fragment else " " + stripped
    return result


def parse_pdf(path: str) -> str:
    logger.debug("Parsing PDF: %s", path)
    pages = []

    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            table_bboxes = [table.bbox for table in page.find_tables()]
            tables = page.extract_tables()
            table_rows = []
            for table in tables:
                for row in table:
                    table_rows.append(" | ".join(_join_cell(cell) if cell else "" for cell in row))
            table_text = "\n".join(table_rows)

            non_table_page = page
            for bbox in table_bboxes:
                non_table_page = non_table_page.filter(lambda obj, bb=bbox: not (bb[0] <= obj["x0"] <= bb[2] and bb[1] <= obj["top"] <= bb[3]))
            words = non_table_page.extract_words(x_tolerance=5, y_tolerance=3)
            words = sorted(words, key=lambda w: (round(w["top"] / 5) * 5, w["x0"]))
            lines = {}
            for w in words:
                line_key = round(w["top"] / 5) * 5
                lines.setdefault(line_key, []).append(w["text"])
            body_text = "\n".join(" ".join(line) for line in lines.values())

            page_text = body_text + "\n" + table_text
            if not page_text.strip():
                logger.warning("No text on page %d, falling back to OCR: %s", page.page_number, path)
                images = convert_from_path(path, first_page=page.page_number, last_page=page.page_number)
                page_text = pytesseract.image_to_string(images[0])
            pages.append(page_text)

    text = "\n".join(pages).strip()
    logger.debug("PDF extracted %d characters", len(text))
    return text


def _cell_str(cell) -> str:
    if isinstance(cell, float):
        if math.isfinite(cell) and cell == int(cell):
            return str(int(cell))
        return f"{cell:.2f}"
    return str(cell)


def parse_xlsx(path: str) -> str:
    logger.debug("Parsing XLSX: %s", path)
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    lines = []
    for sheet in wb.worksheets:
        if getattr(sheet, "sheet_state", "visible") in ("hidden", "veryHidden"):
            continue
        if any(skip in sheet.title for skip in _SKIP_SHEETS):
            continue
        lines.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            str_cells = [_cell_str(c) if c is not None else "" for c in row]
            while str_cells and not str_cells[-1].strip():
                str_cells.pop()
            if any(c.strip() for c in str_cells):
                lines.append(" | ".join(str_cells))
    wb.close()
    text = "\n".join(lines).strip()
    logger.debug("XLSX extracted %d characters", len(text))
    return text


def parse_xls(path: str) -> str:
    try:
        import xlrd
    except ImportError:
        logger.error("xlrd not installed — run: pip install xlrd")
        return ""
    logger.debug("Parsing XLS: %s", path)
    wb = xlrd.open_workbook(path)
    lines = []
    for sheet in wb.sheets():
        if sheet.visibility != 0:
            continue
        if any(skip in sheet.name for skip in _SKIP_SHEETS):
            continue
        lines.append(f"[Sheet: {sheet.name}]")
        for row_idx in range(sheet.nrows):
            str_cells = [_cell_str(sheet.cell_value(row_idx, c)) for c in range(sheet.ncols)]
            while str_cells and not str_cells[-1].strip():
                str_cells.pop()
            if any(c.strip() for c in str_cells):
                lines.append(" | ".join(str_cells))
    text = "\n".join(lines).strip()
    logger.debug("XLS extracted %d characters", len(text))
    return text


def parse_csv(path: str) -> str:
    logger.debug("Parsing CSV: %s", path)
    lines = []
    with open(path, "r", encoding="utf-8-sig", errors="ignore", newline="") as f:
        for row in csv.reader(f):
            str_cells = [c.strip() for c in row]
            while str_cells and not str_cells[-1]:
                str_cells.pop()
            if any(str_cells):
                lines.append(" | ".join(str_cells))
    text = "\n".join(lines).strip()
    logger.debug("CSV extracted %d characters", len(text))
    return text


def parse_spreadsheet(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".xls":
        return parse_xls(path)
    if ext == ".csv":
        return parse_csv(path)
    return parse_xlsx(path)


def _is_old_ppt_format(path: str) -> bool:
    with open(path, "rb") as f:
        return f.read(4) == b"\xD0\xCF\x11\xE0"


def _extract_text_via_com(path: str, powerpoint=None) -> str:
    import win32com.client
    abs_path = os.path.abspath(path)
    owns_instance = powerpoint is None
    if owns_instance:
        import pythoncom
        pythoncom.CoInitialize()
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1
    try:
        deck = powerpoint.Presentations.Open(abs_path, WithWindow=False)
        text = ""
        for i in range(1, deck.Slides.Count + 1):
            text += f"[Slide {i}]\n"
            slide = deck.Slides(i)
            for j in range(1, slide.Shapes.Count + 1):
                shape = slide.Shapes(j)
                try:
                    if shape.HasTextFrame:
                        shape_text = shape.TextFrame.TextRange.Text.strip()
                        if shape_text:
                            text += shape_text + "\n"
                    if shape.HasTable:
                        table = shape.Table
                        for r in range(1, table.Rows.Count + 1):
                            row_cells = [
                                table.Cell(r, c).Shape.TextFrame.TextRange.Text.strip()
                                for c in range(1, table.Columns.Count + 1)
                            ]
                            row_text = " | ".join(row_cells)
                            if row_text.strip("|").strip():
                                text += row_text + "\n"
                    if shape.Type in (11, 13):  # msoLinkedPicture, msoPicture
                        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as _tmp:
                            img_path = _tmp.name
                        try:
                            shape.Export(img_path, 2)  # 2 = ppShapeFormatPNG
                            with Image.open(img_path) as img:
                                ocr_text = pytesseract.image_to_string(img.copy()).strip()
                            if ocr_text:
                                logger.debug("OCR on image in slide %d (old PPT)", i)
                                text += ocr_text + "\n"
                        finally:
                            if os.path.exists(img_path):
                                os.remove(img_path)
                except Exception:
                    pass
        deck.Close()
        logger.debug("COM extracted %d characters from old PPT: %s", len(text), os.path.basename(path))
        return text.strip()
    except Exception as e:
        logger.warning("COM text extraction failed for %s: %s", os.path.basename(path), e)
        return ""
    finally:
        if owns_instance:
            powerpoint.Quit()
            try:
                import pythoncom
                pythoncom.CoUninitialize()
            except Exception:
                pass


def _iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def parse_ppt(path: str, powerpoint=None) -> str:
    logger.debug("Parsing PPT: %s", path)
    if _is_old_ppt_format(path):
        logger.debug("Old binary PPT format, extracting via COM directly: %s", path)
        return _extract_text_via_com(path, powerpoint)
    try:
        prs = Presentation(path)
    except Exception as e:
        logger.warning("Could not open PPTX %s: %s", os.path.basename(path), e)
        return ""
    text = ""
    for slide_num, slide in enumerate(prs.slides, start=1):
        text += f"[Slide {slide_num}]\n"
        for shape in _iter_shapes(slide.shapes):
            if shape.has_text_frame and not (shape.is_placeholder and shape.placeholder_format.idx == 12):
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        text += line + "\n"
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.replace("\n", " ").strip() for cell in row.cells)
                    if row_text.strip("|").strip():
                        text += row_text + "\n"
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = Image.open(io.BytesIO(shape.image.blob))
                ocr_text = pytesseract.image_to_string(image).strip()
                if ocr_text:
                    logger.debug("OCR on image in slide %d", slide_num)
                    text += ocr_text + "\n"
    logger.debug("PPT extracted %d characters", len(text))
    return text.strip()


def parse_zip(path: str, depth: int = 0) -> str:
    import pythoncom
    if depth > 2:
        logger.warning("Skipping deeply nested ZIP at depth %d: %s", depth, path)
        return ""
    logger.debug("Parsing ZIP: %s", path)
    text = ""
    temp_dir = tempfile.mkdtemp()
    powerpoint = None
    com_initialized = False
    try:
        with zipfile.ZipFile(path, "r") as z:
            z.extractall(temp_dir)
        all_files = [
            (os.path.join(root, filename), filename)
            for root, _, files in os.walk(temp_dir)
            for filename in files
        ]
        for file_path, filename in all_files:
            ext = os.path.splitext(filename)[-1].lower()
            logger.debug("ZIP entry: %s", filename)
            if ext in PDF_EXTENSIONS:
                text += f"[{filename}]\n" + parse_pdf(file_path) + "\n"
            elif ext in SPREADSHEET_EXTENSIONS:
                text += f"[{filename}]\n" + parse_spreadsheet(file_path) + "\n"
            elif ext in PRESENTATION_EXTENSIONS:
                if powerpoint is None and _is_old_ppt_format(file_path):
                    try:
                        import win32com.client
                        pythoncom.CoInitialize()
                        com_initialized = True
                        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                        powerpoint.Visible = 1
                    except Exception as e:
                        logger.warning("Could not start PowerPoint for COM conversion: %s", e)
                text += f"[{filename}]\n" + parse_ppt(file_path, powerpoint=powerpoint) + "\n"
            elif ext in DOCUMENT_EXTENSIONS:
                text += f"[{filename}]\n" + parse_txt(file_path) + "\n"
            elif ext in ZIP_EXTENSIONS:
                text += f"[{filename}]\n" + parse_zip(file_path, depth=depth + 1) + "\n"
            else:
                logger.debug("Skipping unsupported file in ZIP: %s", filename)
    finally:
        if powerpoint is not None:
            powerpoint.Quit()
        if com_initialized:
            pythoncom.CoUninitialize()
        shutil.rmtree(temp_dir, ignore_errors=True)
    logger.debug("ZIP extracted %d characters", len(text))
    return text.strip()


def parse_txt(path: str) -> str:
    logger.debug("Parsing TXT: %s", path)
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read().strip()
    logger.debug("TXT extracted %d characters", len(text))
    return text


def parse_single_attachment(path: str, filename: str) -> tuple:
    ext = os.path.splitext(filename)[-1].lower()
    logger.debug("Dispatching: %s (%s)", filename, ext)

    if ext in PDF_EXTENSIONS:
        text = parse_pdf(path)
    elif ext in SPREADSHEET_EXTENSIONS:
        text = parse_spreadsheet(path)
    elif ext in PRESENTATION_EXTENSIONS:
        text = parse_ppt(path)
    elif ext in DOCUMENT_EXTENSIONS:
        text = parse_txt(path)
    elif ext in ZIP_EXTENSIONS:
        text = parse_zip(path)
    elif ext in IMAGE_EXTENSIONS:
        image = Image.open(path)
        text = pytesseract.image_to_string(image).strip()
        logger.debug("OCR on image %s: %d characters", filename, len(text))
    else:
        logger.warning("Unsupported file type %s, skipping: %s", ext, filename)
        text = ""

    return (filename, text)
